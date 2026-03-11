#!/usr/bin/env python3
"""
VFT Document Processor — Text Extraction Layer

Extracts plain text from dataroom files (PDFs, spreadsheets, presentations)
and stores per-page content in the document_pages table. This is the
"symbolic handle" layer — no LLM calls, pure Python extraction.

Usage:
    python extract_text.py /path/to/dataroom --slug midbound_dataroom
    python extract_text.py /path/to/single/file.pdf --slug midbound_dataroom
    python extract_text.py /path/to/dataroom --slug midbound_dataroom --dry-run
"""

import argparse
import json
import os
import sqlite3
import sys
import tempfile
import zipfile
from datetime import datetime, UTC
from pathlib import Path

REPO_ROOT = Path(__file__).resolve().parents[3]
DB_PATH = REPO_ROOT / "fund" / "metadata" / "db" / "ingestion.db"

# Supported file extensions mapped to extraction method
EXTRACTORS = {
    ".pdf": "pdf",
    ".xlsx": "xlsx",
    ".xls": "xlsx",
    ".pptx": "pptx",
    ".ppt": "pptx",
    ".csv": "text",
    ".txt": "text",
    ".md": "text",
    ".doc": "text",      # Limited support — plain text only
    ".docx": "text",     # Limited support — plain text only
    ".json": "text",
    ".zip": "zip",
}

# Pages with fewer chars than this after native PDF extraction trigger OCR
OCR_CHAR_THRESHOLD = 50


def get_db() -> sqlite3.Connection:
    conn = sqlite3.connect(str(DB_PATH))
    conn.execute("PRAGMA journal_mode=OFF")
    conn.execute("PRAGMA foreign_keys=ON")
    conn.row_factory = sqlite3.Row
    return conn


# ── PDF Extraction ──────────────────────────────────────────────────────

def extract_pdf_native(file_path: Path) -> list[dict]:
    """Extract text from PDF using pymupdf (native text layer).

    Returns list of page dicts. Pages with < OCR_CHAR_THRESHOLD chars
    are flagged for potential OCR fallback.
    """
    try:
        import fitz  # pymupdf
    except ImportError:
        print("[VFT] pymupdf not installed. Run: pip install pymupdf", file=sys.stderr)
        return []

    pages = []
    try:
        doc = fitz.open(str(file_path))
        total = len(doc)
        for i, page in enumerate(doc):
            text = page.get_text("text")
            char_count = len(text.strip())
            has_images = len(page.get_images()) > 0

            pages.append({
                "page": i + 1,
                "total_pages": total,
                "text": text,
                "char_count": char_count,
                "method": "native_pdf",
                "quality": 1.0 if char_count >= OCR_CHAR_THRESHOLD else 0.3,
                "needs_ocr": char_count < OCR_CHAR_THRESHOLD and has_images,
            })
        doc.close()
    except Exception as e:
        print(f"[VFT] PDF native extraction failed for {file_path}: {e}", file=sys.stderr)

    return pages


def extract_pdf_ocr(file_path: Path, page_numbers: list[int] = None) -> list[dict]:
    """OCR specific pages of a PDF using pytesseract.

    Only runs on pages that need it (scanned/image-based).
    """
    try:
        import fitz  # pymupdf
        import pytesseract
        from PIL import Image
        import io
    except ImportError as e:
        print(f"[VFT] OCR dependencies missing ({e}). Run: pip install pymupdf pytesseract Pillow", file=sys.stderr)
        return []

    pages = []
    try:
        doc = fitz.open(str(file_path))
        total = len(doc)
        target_pages = page_numbers or range(total)

        for i in target_pages:
            page = doc[i]
            # Render page to image at 300 DPI for OCR
            mat = fitz.Matrix(300 / 72, 300 / 72)
            pix = page.get_pixmap(matrix=mat)
            img_bytes = pix.tobytes("png")
            img = Image.open(io.BytesIO(img_bytes))

            text = pytesseract.image_to_string(img)
            # Get OCR confidence
            ocr_data = pytesseract.image_to_data(img, output_type=pytesseract.Output.DICT)
            confidences = [int(c) for c in ocr_data["conf"] if str(c).isdigit() and int(c) > 0]
            avg_confidence = sum(confidences) / len(confidences) / 100.0 if confidences else 0.5

            pages.append({
                "page": i + 1,
                "total_pages": total,
                "text": text,
                "char_count": len(text.strip()),
                "method": "ocr",
                "quality": round(avg_confidence, 3),
                "needs_ocr": False,
            })
        doc.close()
    except Exception as e:
        print(f"[VFT] OCR extraction failed for {file_path}: {e}", file=sys.stderr)

    return pages


def extract_pdf(file_path: Path) -> list[dict]:
    """Full PDF extraction: native first, OCR fallback for scanned pages."""
    pages = extract_pdf_native(file_path)
    if not pages:
        return []

    # Identify pages needing OCR
    ocr_needed = [p["page"] - 1 for p in pages if p.get("needs_ocr")]
    if ocr_needed:
        print(f"[VFT] {len(ocr_needed)} pages need OCR in {file_path.name}")
        ocr_pages = extract_pdf_ocr(file_path, ocr_needed)
        # Replace native results with OCR results for those pages
        ocr_map = {p["page"]: p for p in ocr_pages}
        for i, page in enumerate(pages):
            if page["page"] in ocr_map:
                pages[i] = ocr_map[page["page"]]

    return pages


# ── Spreadsheet Extraction ──────────────────────────────────────────────

def extract_xlsx(file_path: Path) -> list[dict]:
    """Extract spreadsheet content, one 'page' per sheet.

    Renders content as markdown tables for LLM readability.
    """
    try:
        import openpyxl
    except ImportError:
        print("[VFT] openpyxl not installed. Run: pip install openpyxl", file=sys.stderr)
        return []

    pages = []
    try:
        wb = openpyxl.load_workbook(str(file_path), read_only=True, data_only=True)
        total_sheets = len(wb.sheetnames)

        for sheet_idx, sheet_name in enumerate(wb.sheetnames):
            ws = wb[sheet_name]
            rows = list(ws.iter_rows(values_only=True))
            if not rows:
                pages.append({
                    "page": sheet_idx + 1,
                    "total_pages": total_sheets,
                    "text": f"# Sheet: {sheet_name}\n\n(empty)",
                    "char_count": 0,
                    "method": "xlsx",
                    "quality": 1.0,
                    "metadata": {"sheet_name": sheet_name},
                })
                continue

            # Build markdown table
            lines = [f"# Sheet: {sheet_name}", ""]

            # Header row
            headers = [str(c) if c is not None else "" for c in rows[0]]
            lines.append("| " + " | ".join(headers) + " |")
            lines.append("| " + " | ".join(["---"] * len(headers)) + " |")

            # Data rows (cap at 500 rows to avoid massive text)
            for row in rows[1:501]:
                cells = [str(c) if c is not None else "" for c in row]
                # Pad or truncate to match header count
                while len(cells) < len(headers):
                    cells.append("")
                cells = cells[:len(headers)]
                lines.append("| " + " | ".join(cells) + " |")

            if len(rows) > 501:
                lines.append(f"\n... ({len(rows) - 501} more rows truncated)")

            text = "\n".join(lines)
            pages.append({
                "page": sheet_idx + 1,
                "total_pages": total_sheets,
                "text": text,
                "char_count": len(text),
                "method": "xlsx",
                "quality": 1.0,
                "metadata": {"sheet_name": sheet_name, "row_count": len(rows)},
            })

        wb.close()
    except Exception as e:
        print(f"[VFT] XLSX extraction failed for {file_path}: {e}", file=sys.stderr)

    return pages


# ── Presentation Extraction ─────────────────────────────────────────────

def extract_pptx(file_path: Path) -> list[dict]:
    """Extract text from presentations, one 'page' per slide."""
    try:
        from pptx import Presentation
    except ImportError:
        print("[VFT] python-pptx not installed. Run: pip install python-pptx", file=sys.stderr)
        return []

    pages = []
    try:
        prs = Presentation(str(file_path))
        total_slides = len(prs.slides)

        for slide_idx, slide in enumerate(prs.slides):
            texts = []
            # Slide title
            if slide.shapes.title:
                texts.append(f"# {slide.shapes.title.text}")
                texts.append("")

            # All text frames
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            texts.append(text)

            # Tables in slide
            if hasattr(shape, "has_table") and shape.has_table:
                table = shape.table
                for row in table.rows:
                    cells = [cell.text for cell in row.cells]
                    texts.append("| " + " | ".join(cells) + " |")

            # Notes
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    texts.append(f"\n---\nSpeaker notes: {notes}")

            full_text = "\n".join(texts)
            pages.append({
                "page": slide_idx + 1,
                "total_pages": total_slides,
                "text": full_text,
                "char_count": len(full_text),
                "method": "pptx",
                "quality": 1.0,
                "metadata": {"slide_number": slide_idx + 1},
            })
    except Exception as e:
        print(f"[VFT] PPTX extraction failed for {file_path}: {e}", file=sys.stderr)

    return pages


# ── Plain Text Extraction ───────────────────────────────────────────────

def extract_text_file(file_path: Path) -> list[dict]:
    """Extract content from plain text files (txt, md, csv, json)."""
    try:
        text = file_path.read_text(encoding="utf-8", errors="replace")
    except Exception as e:
        print(f"[VFT] Text extraction failed for {file_path}: {e}", file=sys.stderr)
        return []

    return [{
        "page": 1,
        "total_pages": 1,
        "text": text,
        "char_count": len(text),
        "method": "text",
        "quality": 1.0,
    }]


# ── ZIP Extraction ──────────────────────────────────────────────────────

def extract_zip(file_path: Path, dataroom_slug: str, conn: sqlite3.Connection) -> dict:
    """Extract ZIP contents to temp dir and recursively process.

    Returns summary dict with counts.
    """
    results = {"extracted": 0, "skipped": 0, "errors": 0}

    try:
        with tempfile.TemporaryDirectory() as tmpdir:
            with zipfile.ZipFile(str(file_path), "r") as zf:
                zf.extractall(tmpdir)

            tmp_path = Path(tmpdir)
            for child in sorted(tmp_path.rglob("*")):
                if not child.is_file():
                    continue
                # Skip macOS resource forks and metadata
                rel = child.relative_to(tmp_path)
                if "__MACOSX" in str(rel) or child.name.startswith("._"):
                    continue
                virtual_path = f"{file_path.name}/{rel}"
                result = extract_file(child, dataroom_slug, conn, override_path=virtual_path)
                if result["status"] == "extracted":
                    results["extracted"] += 1
                elif result["status"] == "skipped":
                    results["skipped"] += 1
                else:
                    results["errors"] += 1
    except Exception as e:
        print(f"[VFT] ZIP extraction failed for {file_path}: {e}", file=sys.stderr)
        results["errors"] += 1

    return results


# ── Main Dispatcher ─────────────────────────────────────────────────────

def extract_file(
    file_path: Path,
    dataroom_slug: str,
    conn: sqlite3.Connection,
    override_path: str = None,
) -> dict:
    """Extract text from a single file and store in document_pages.

    Returns: {"status": "extracted"|"skipped"|"error"|"unsupported", "pages": int, "method": str}
    """
    ext = file_path.suffix.lower()
    extractor_type = EXTRACTORS.get(ext)

    if not extractor_type:
        return {"status": "unsupported", "pages": 0, "method": None}

    # Use override_path for files extracted from ZIPs
    db_path = override_path or str(file_path)

    # Check if already extracted (dedup)
    existing = conn.execute(
        "SELECT COUNT(*) FROM document_pages WHERE file_path = ? AND dataroom_slug = ?",
        (db_path, dataroom_slug),
    ).fetchone()[0]
    if existing > 0:
        return {"status": "skipped", "pages": existing, "method": "cached"}

    # Handle ZIP separately (recursive)
    if extractor_type == "zip":
        zip_result = extract_zip(file_path, dataroom_slug, conn)
        return {"status": "extracted", "pages": zip_result["extracted"], "method": "zip"}

    # Extract pages
    if extractor_type == "pdf":
        pages = extract_pdf(file_path)
    elif extractor_type == "xlsx":
        pages = extract_xlsx(file_path)
    elif extractor_type == "pptx":
        pages = extract_pptx(file_path)
    elif extractor_type == "text":
        pages = extract_text_file(file_path)
    else:
        return {"status": "unsupported", "pages": 0, "method": None}

    if not pages:
        return {"status": "error", "pages": 0, "method": extractor_type}

    # Store in document_pages
    for page in pages:
        metadata = json.dumps(page.get("metadata", {}))
        conn.execute(
            """INSERT OR IGNORE INTO document_pages
               (file_path, dataroom_slug, page_number, total_pages,
                text_content, char_count, extraction_method, extraction_quality, metadata)
               VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?)""",
            (
                db_path,
                dataroom_slug,
                page["page"],
                page["total_pages"],
                page["text"],
                page["char_count"],
                page["method"],
                page.get("quality", 1.0),
                metadata,
            ),
        )

    conn.commit()
    return {
        "status": "extracted",
        "pages": len(pages),
        "method": pages[0]["method"] if pages else extractor_type,
    }


def extract_directory(
    root: Path,
    dataroom_slug: str,
    conn: sqlite3.Connection,
    dry_run: bool = False,
) -> dict:
    """Extract text from all supported files in a directory.

    Returns summary dict.
    """
    summary = {"extracted": 0, "skipped": 0, "unsupported": 0, "errors": 0, "total_pages": 0, "files": []}

    for file_path in sorted(root.rglob("*")):
        if not file_path.is_file():
            continue
        # Skip hidden files and macOS resource forks
        if file_path.name.startswith(".") or file_path.name.startswith("__"):
            continue

        rel_path = str(file_path.relative_to(root))

        if dry_run:
            ext = file_path.suffix.lower()
            supported = ext in EXTRACTORS
            summary["files"].append({
                "path": rel_path,
                "supported": supported,
                "method": EXTRACTORS.get(ext, "unsupported"),
                "size_bytes": file_path.stat().st_size,
            })
            if supported:
                summary["extracted"] += 1
            else:
                summary["unsupported"] += 1
            continue

        result = extract_file(file_path, dataroom_slug, conn, override_path=rel_path)
        summary["files"].append({"path": rel_path, **result})

        if result["status"] == "extracted":
            summary["extracted"] += 1
            summary["total_pages"] += result["pages"]
        elif result["status"] == "skipped":
            summary["skipped"] += 1
            summary["total_pages"] += result["pages"]
        elif result["status"] == "unsupported":
            summary["unsupported"] += 1
        else:
            summary["errors"] += 1

    return summary


if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="VFT Document Text Extraction")
    parser.add_argument("path", help="File or directory to extract")
    parser.add_argument("--slug", required=True, help="Dataroom slug identifier")
    parser.add_argument("--db", type=str, help="Path to ingestion.db")
    parser.add_argument("--dry-run", action="store_true", help="Preview without extracting")
    args = parser.parse_args()

    db_path = args.db or str(DB_PATH)
    if not os.path.exists(db_path):
        print(f"[VFT] Database not found at {db_path}. Run init_db.py first.", file=sys.stderr)
        sys.exit(1)

    conn = get_db() if not args.db else sqlite3.connect(args.db)
    if args.db:
        conn.execute("PRAGMA journal_mode=OFF")
        conn.row_factory = sqlite3.Row

    target = Path(args.path).expanduser().resolve()

    if target.is_dir():
        result = extract_directory(target, args.slug, conn, dry_run=args.dry_run)
        print(json.dumps({
            "dataroom": args.slug,
            "extracted": result["extracted"],
            "skipped": result["skipped"],
            "unsupported": result["unsupported"],
            "errors": result["errors"],
            "total_pages": result["total_pages"],
        }, indent=2))
    elif target.is_file():
        result = extract_file(target, args.slug, conn)
        print(json.dumps(result, indent=2))
    else:
        print(f"[VFT] Path not found: {target}", file=sys.stderr)
        sys.exit(1)

    conn.close()
